import glob
import io
import os
import re
import sys
import configparser
import tkinter as tk
from tkinter import filedialog, messagebox
from datetime import date, timedelta

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.oxml.ns import qn

CONFIG_FILENAME = 'config.ini'

BG       = '#0f1428'
CARD_BG  = '#1a2040'
ACCENT   = '#5080e0'
FG       = '#ffffff'
FG_DIM   = '#8898cc'
BORDER   = '#2a3560'
SUCCESS  = '#2a6040'
ERROR    = '#602a40'
FONT     = '맑은 고딕'


def _base_dir() -> str:
    return os.path.dirname(sys.executable if getattr(sys, 'frozen', False) else os.path.abspath(__file__))


def _config_path() -> str:
    return os.path.join(_base_dir(), CONFIG_FILENAME)


def load_config() -> tuple[str, str]:
    cfg = configparser.ConfigParser()
    cfg.read(_config_path(), encoding='utf-8')
    views = cfg.get('paths', 'views_dir', fallback='').strip()
    pptx  = cfg.get('paths', 'pptx_dir',  fallback='').strip()
    return views, pptx


def save_config(views: str, pptx: str) -> None:
    cfg = configparser.ConfigParser()
    cfg['paths'] = {'views_dir': views, 'pptx_dir': pptx}
    with open(_config_path(), 'w', encoding='utf-8') as f:
        cfg.write(f)


def find_numbered_images(folder: str) -> list[str]:
    import re
    if not os.path.isdir(folder):
        return []
    files = [
        f for f in os.listdir(folder)
        if os.path.splitext(f)[1].lower() in ('.jpg', '.jpeg', '.png', '.gif', '.bmp')
    ]
    files.sort(key=lambda f: [int(t) if t.isdigit() else t.lower() for t in re.split(r'(\d+)', f)])
    return [os.path.join(folder, f) for f in files]


def next_wednesday() -> str:
    today = date.today()
    days = (2 - today.weekday()) % 7 or 7
    return (today + timedelta(days=days)).strftime('%Y%m%d')


_PPTX_SUPPORTED_FORMATS = {'BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF'}


def _normalize_image(raw: bytes) -> io.BytesIO:
    """python-pptx가 지원하지 않는 포맷(MPO 등)을 JPEG로 변환."""
    from PIL import Image  # noqa: PLC0415

    img = Image.open(io.BytesIO(raw))
    fmt = (img.format or '').upper()
    if fmt in _PPTX_SUPPORTED_FORMATS:
        return io.BytesIO(raw)

    buf = io.BytesIO()
    img.convert('RGB').save(buf, format='JPEG', quality=92)
    buf.seek(0)
    return buf


def parse_pages(pattern: str, total: int) -> list[int]:
    """
    페이지 구성 패턴 파싱.
      "2"        → [2, 2, ...] (전체 이미지를 2장씩 나눠 담음)
      "2,2,3"    → [2, 2, 3]
      "2x3"      → [2, 2, 2]
      "2x2,3x1"  → [2, 2, 3]
    """
    pattern = (pattern or '').strip()
    if not pattern:
        return [total] if total else []

    chunks = [c.strip() for c in re.split(r'[,+]', pattern) if c.strip()]
    is_single_number = len(chunks) == 1 and chunks[0].isdigit()

    result: list[int] = []
    for chunk in chunks:
        m = re.match(r'^(\d+)\s*[x*X]\s*(\d+)$', chunk)
        if m:
            n, rep = int(m.group(1)), int(m.group(2))
            if n <= 0 or rep <= 0:
                raise ValueError(f'잘못된 패턴: {chunk!r}')
            result.extend([n] * rep)
        elif chunk.isdigit():
            n = int(chunk)
            if n <= 0:
                raise ValueError(f'잘못된 패턴: {chunk!r}')
            result.append(n)
        else:
            raise ValueError(f'잘못된 패턴: {chunk!r}')

    if is_single_number and result:
        per = result[0]
        result = []
        remaining = total
        while remaining > 0:
            take = min(per, remaining)
            result.append(take)
            remaining -= take

    return result


def find_or_create_pptx(folder: str, title: str) -> tuple[str, str | None]:
    """(저장할 경로, 기존 파일 경로 or None) 반환. 항상 {title}.pptx로 저장."""
    os.makedirs(folder, exist_ok=True)
    save_path = os.path.join(folder, f'{title}.pptx')
    files = [f for f in glob.glob(os.path.join(folder, '*.pptx'))
             if not os.path.basename(f).startswith('~$')]
    existing = files[0] if files else None
    if not existing:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(save_path)
    return save_path, existing


def run_update(views_dir: str, pptx_dir: str, title: str, pattern: str) -> str:
    images = find_numbered_images(views_dir)
    if not images:
        raise ValueError(f'이미지가 없습니다.\n폴더: {views_dir}')

    total = len(images)
    sizes = parse_pages(pattern, total)
    if not sizes:
        raise ValueError('페이지 구성이 비어 있습니다.')

    groups: list[list[str]] = []
    idx = 0
    for size in sizes:
        if idx >= total:
            break
        groups.append(images[idx:idx + size])
        idx += size
    if idx < total:
        groups.append(images[idx:])

    save_path, existing = find_or_create_pptx(pptx_dir, title)
    prs = Presentation(existing or save_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 슬라이드 수 맞추기
    blank = prs.slide_layouts[6]
    needed = len([g for g in groups if g])
    while len(prs.slides) < needed:
        prs.slides.add_slide(blank)
    while len(prs.slides) > needed:
        sld_lst = list(prs.slides._sldIdLst)
        last = sld_lst[-1]
        rId = last.get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(last)

    slide_num = 0
    for group in groups:
        if not group:
            continue
        slide = prs.slides[slide_num]
        sp_tree = slide.shapes._spTree
        for shape in list(slide.shapes):
            sp_tree.remove(shape._element)

        top_offset = Emu(0)
        if slide_num == 0:
            m = re.match(r'^(\d{4})(\d{2})(\d{2})\s+(.+?)(?:\s+콘티)?$', title)
            if m:
                formatted_date = f'{m.group(1)}.{int(m.group(2))}.{int(m.group(3))}'
                label = m.group(4)
            else:
                formatted_date = ''
                label = title

            NAVY = RGBColor(0x1A, 0x2A, 0x4A)
            WHITE = RGBColor(0xFF, 0xFF, 0xFF)
            box_h = Pt(32)
            label_w = Pt(28 + len(label) * 18)

            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), label_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = NAVY
            box.line.fill.background()
            tf = box.text_frame
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.name = '맑은 고딕'
            run.font.color.rgb = WHITE

            if formatted_date:
                date_box = slide.shapes.add_textbox(label_w + Pt(10), Emu(0), Pt(180), box_h)
                dtf = date_box.text_frame
                dtf.margin_top = Emu(0)
                dtf.margin_bottom = Emu(0)
                dtf.margin_left = Emu(0)
                dtf.margin_right = Emu(0)
                dtf.word_wrap = False
                dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
                dp = dtf.paragraphs[0]
                dp.alignment = PP_ALIGN.LEFT
                drun = dp.add_run()
                drun.text = formatted_date
                drun.font.size = Pt(18)
                drun.font.bold = True
                drun.font.name = '맑은 고딕'
                drun.font.color.rgb = NAVY

            top_offset = box_h + Pt(4)

        n = len(group)
        img_w = slide_w // n
        img_h = slide_h - top_offset
        for i, img_path in enumerate(group):
            with open(img_path, 'rb') as f:
                img_data = _normalize_image(f.read())
            slide.shapes.add_picture(img_data, img_w * i, top_offset, img_w, img_h)

        slide_num += 1

    prs.save(save_path)
    return f'({total}장 → {needed}슬라이드) {os.path.basename(save_path)}', save_path


class CollapsiblePanel(tk.Frame):
    """클릭하면 부드럽게 펼쳐지는 토글 패널. 펼침 시 항목이 한 줄씩 슬라이드 인."""

    HEADER_BG = '#222b4a'
    HEADER_HOVER = '#2b375c'
    BODY_BG = '#162038'
    ACCENT_DIM = '#6f86c2'

    def __init__(self, parent, title: str, lines: list[tuple[str, str]]):
        super().__init__(parent, bg=CARD_BG)
        self.expanded = False
        self._title = title
        self._lines = lines
        self._anim_after = None

        self.header = tk.Frame(self, bg=self.HEADER_BG, cursor='hand2', height=30)
        self.header.pack(fill='x')
        self.header.pack_propagate(False)

        inner = tk.Frame(self.header, bg=self.HEADER_BG, cursor='hand2')
        inner.place(relx=0, rely=0.5, anchor='w', x=10)

        self.chevron = tk.Label(inner, text='▸', bg=self.HEADER_BG, fg=ACCENT,
                                font=(FONT, 11, 'bold'), cursor='hand2')
        self.chevron.pack(side='left', padx=(0, 8))

        self.title_label = tk.Label(inner, text=title, bg=self.HEADER_BG, fg=FG,
                                     font=(FONT, 10, 'bold'), cursor='hand2')
        self.title_label.pack(side='left')

        self.hint = tk.Label(self.header, text='클릭', bg=self.HEADER_BG,
                              fg=self.ACCENT_DIM, font=(FONT, 8), cursor='hand2')
        self.hint.place(relx=1.0, rely=0.5, anchor='e', x=-12)

        for w in (self.header, inner, self.chevron, self.title_label, self.hint):
            w.bind('<Button-1>', lambda e: self.toggle())
            w.bind('<Enter>', self._on_enter)
            w.bind('<Leave>', self._on_leave)

        self.body = tk.Frame(self, bg=self.BODY_BG)
        self._line_widgets: list[tk.Frame] = []
        for marker, text in lines:
            row = tk.Frame(self.body, bg=self.BODY_BG)
            mark_color = ACCENT if marker.isdigit() else self.ACCENT_DIM
            tk.Label(row, text=marker, bg=self.BODY_BG, fg=mark_color,
                     font=(FONT, 9, 'bold'), width=2, anchor='center').pack(side='left', padx=(12, 6))
            tk.Label(row, text=text, bg=self.BODY_BG, fg=FG,
                     font=(FONT, 9), anchor='w', justify='left',
                     wraplength=320).pack(side='left', fill='x', expand=True, pady=1)
            self._line_widgets.append(row)

    def _on_enter(self, _):
        for w in (self.header, self.chevron, self.title_label, self.hint):
            w.configure(bg=self.HEADER_HOVER)

    def _on_leave(self, _):
        for w in (self.header, self.chevron, self.title_label, self.hint):
            w.configure(bg=self.HEADER_BG)

    def toggle(self):
        if self.expanded:
            self._collapse()
        else:
            self._expand()

    def _expand(self):
        self.expanded = True
        self.hint.configure(text='닫기')
        self.body.pack(fill='x', pady=(0, 0))
        for w in self._line_widgets:
            w.pack_forget()
        self._reveal_idx = 0
        self._spin_chevron(0, opening=True)
        self._reveal_next()

    def _collapse(self):
        self.expanded = False
        self.hint.configure(text='클릭')
        if self._anim_after:
            try:
                self.after_cancel(self._anim_after)
            except Exception:
                pass
            self._anim_after = None
        self.body.pack_forget()
        self._spin_chevron(0, opening=False)

    def _reveal_next(self):
        if self._reveal_idx >= len(self._line_widgets):
            self._anim_after = None
            return
        self._line_widgets[self._reveal_idx].pack(fill='x', pady=1)
        self._reveal_idx += 1
        self._anim_after = self.after(35, self._reveal_next)

    def _spin_chevron(self, step: int, opening: bool):
        frames_open  = ['▸', '▹', '▾', '▼']
        frames_close = ['▼', '▿', '▹', '▸']
        frames = frames_open if opening else frames_close
        if step >= len(frames):
            return
        self.chevron.configure(text=frames[step])
        self.after(40, lambda: self._spin_chevron(step + 1, opening))


class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title('Setlist')
        self.resizable(False, False)
        self.configure(bg=BG)

        views, pptx = load_config()
        self.views_var = tk.StringVar(value=views)
        self.pptx_var  = tk.StringVar(value=pptx)

        wed = next_wednesday()
        self.title_var = tk.StringVar(value=f'{wed} 수요성령집회 콘티')
        self.pattern_var = tk.StringVar(value='2')
        self._last_pptx = None

        self._build_ui()
        self._center()
        self.bind('<Return>', lambda e: self._run())
        self.bind('<KP_Enter>', lambda e: self._run())

    def _label(self, parent, text):
        return tk.Label(parent, text=text, bg=CARD_BG, fg=FG_DIM, font=(FONT, 9))

    def _entry(self, parent, var, width=38):
        return tk.Entry(parent, textvariable=var, width=width,
                        bg=BG, fg=FG, insertbackground=FG,
                        relief='flat', highlightthickness=1,
                        highlightbackground=BORDER, highlightcolor=ACCENT,
                        font=(FONT, 10))

    def _btn(self, parent, text, cmd, small=False):
        return tk.Button(parent, text=text, command=cmd,
                         bg=ACCENT, fg=FG, activebackground='#3a60c0',
                         activeforeground=FG, relief='flat', cursor='hand2',
                         font=(FONT, 9 if small else 11, 'bold'),
                         padx=8 if small else 16, pady=4 if small else 10)

    def _folder_row(self, parent, label, var, row):
        self._label(parent, label).grid(row=row, column=0, columnspan=2, sticky='w', pady=(12, 2))
        e = self._entry(parent, var, width=34)
        e.grid(row=row+1, column=0, sticky='ew', padx=(0, 6))
        self._btn(parent, '변경', lambda: self._pick(var), small=True).grid(row=row+1, column=1)

    def _build_ui(self):
        card = tk.Frame(self, bg=CARD_BG, padx=28, pady=24)
        card.pack(padx=16, pady=16)

        tk.Label(card, text='Setlist', bg=CARD_BG, fg=FG,
                 font=(FONT, 16, 'bold')).grid(row=0, column=0, columnspan=2, sticky='w', pady=(0, 12))

        help_panel = CollapsiblePanel(
            card,
            title='사용법',
            lines=[
                ('1', '이미지 폴더에 1.jpg, 2.jpg 순서로 악보 이미지 저장'),
                ('2', '이미지 폴더 지정'),
                ('3', '제목을 "YYYYMMDD 라벨" 형식으로 (예: 20260520 수요성령집회)'),
                ('4', '페이지 구성 입력 — 아래 예시 참고'),
                ('•', '2  →  2장씩 모두 채움'),
                ('•', '2,2,3  →  3슬라이드 (2장, 2장, 3장)'),
                ('•', '2x3  →  2장씩 3슬라이드'),
                ('5', '실행 클릭'),
            ],
        )
        help_panel.grid(row=2, column=0, columnspan=2, sticky='ew', pady=(0, 4))

        self._folder_row(card, '이미지 폴더 (views)', self.views_var, 3)
        self._folder_row(card, 'PPT 저장 폴더', self.pptx_var, 5)

        self._label(card, '제목').grid(row=7, column=0, columnspan=2, sticky='w', pady=(12, 2))
        self._entry(card, self.title_var).grid(row=8, column=0, columnspan=2, sticky='ew')

        self._label(card, '페이지 구성  (예: 2 / 2,2,3 / 2x3)').grid(row=9, column=0, columnspan=2, sticky='w', pady=(12, 2))
        self._entry(card, self.pattern_var, width=16).grid(row=10, column=0, sticky='w')

        self._btn(card, '실행', self._run).grid(row=11, column=0, columnspan=2, sticky='ew', pady=(20, 0))

        self.status = tk.Label(card, text='', bg=CARD_BG, fg=FG_DIM,
                               font=(FONT, 9), wraplength=360, justify='left')
        self.status.grid(row=12, column=0, columnspan=2, sticky='w', pady=(10, 0))

        btn_frame = tk.Frame(card, bg=CARD_BG)
        btn_frame.grid(row=13, column=0, columnspan=2, sticky='ew', pady=(8, 0))
        btn_frame.columnconfigure(0, weight=1)
        btn_frame.columnconfigure(1, weight=1)

        self.open_file_btn = self._btn(btn_frame, '파일 열기', self._open_file_direct)
        self.open_file_btn.grid(row=0, column=0, sticky='ew', padx=(0, 4))
        self.open_file_btn.grid_remove()

        self.open_folder_btn = self._btn(btn_frame, '폴더 열기', self._open_folder)
        self.open_folder_btn.grid(row=0, column=1, sticky='ew', padx=(4, 0))
        self.open_folder_btn.grid_remove()

        card.columnconfigure(0, weight=1)

    def _pick(self, var):
        path = filedialog.askdirectory(initialdir=var.get() or os.path.expanduser('~'))
        if path:
            var.set(path)
            save_config(self.views_var.get(), self.pptx_var.get())

    def _open_folder(self):
        if self._last_pptx and os.path.exists(self._last_pptx):
            os.startfile(os.path.dirname(self._last_pptx))

    def _open_file_direct(self):
        if self._last_pptx and os.path.exists(self._last_pptx):
            os.startfile(self._last_pptx)

    def _run(self):
        views = self.views_var.get().strip()
        pptx  = self.pptx_var.get().strip()
        title = self.title_var.get().strip()
        pattern = self.pattern_var.get().strip()

        if not views or not pptx:
            messagebox.showwarning('경고', '폴더를 선택해주세요.')
            return
        if not title:
            messagebox.showwarning('경고', '제목을 입력해주세요.')
            return
        if not pattern:
            messagebox.showwarning('경고', '페이지 구성을 입력해주세요.')
            return

        save_config(views, pptx)
        self.open_folder_btn.grid_remove()
        self.open_file_btn.grid_remove()
        self.status.config(text='실행 중...', fg=FG_DIM)
        self.update()

        try:
            msg, pptx_path = run_update(views, pptx, title, pattern)
            self._last_pptx = pptx_path
            self.status.config(text=f'완료! {msg}', fg='#80e0a0')
            self.open_folder_btn.grid()
            self.open_file_btn.grid()
        except PermissionError:
            self.status.config(text='❌ PPTX 파일이 열려있습니다. PowerPoint를 닫고 다시 실행하세요.', fg='#e08080')
        except Exception as e:
            self.status.config(text=f'❌ {e}', fg='#e08080')

    def _center(self):
        self.update_idletasks()
        w, h = self.winfo_width(), self.winfo_height()
        x = (self.winfo_screenwidth()  - w) // 2
        y = (self.winfo_screenheight() - h) // 2
        self.geometry(f'+{x}+{y}')


if __name__ == '__main__':
    App().mainloop()
