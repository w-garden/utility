import io
import json
import os
import re
import uuid
from datetime import date, datetime, timedelta, timezone
from urllib.parse import quote

import boto3

BUCKET = os.environ["BUCKET_NAME"]
APP_PASSWORD = os.environ.get("APP_PASSWORD", "")
LOG_TABLE = os.environ.get("LOG_TABLE", "")
s3 = boto3.client("s3", region_name="ap-northeast-2",
                  endpoint_url="https://s3.ap-northeast-2.amazonaws.com")
dynamodb = boto3.resource("dynamodb", region_name="ap-northeast-2")


def _log(event: dict, action: str, extra: dict = None):
    if not LOG_TABLE:
        return
    try:
        now = datetime.now(timezone.utc)
        ip = (event.get("requestContext", {})
                   .get("identity", {})
                   .get("sourceIp", "unknown"))
        item = {
            "id": str(uuid.uuid4()),
            "timestamp": now.isoformat(),
            "action": action,
            "ip": ip,
        }
        if extra:
            item.update(extra)
        dynamodb.Table(LOG_TABLE).put_item(Item=item)
    except Exception:
        pass


def _html(body: str, status: int = 200) -> dict:
    return {"statusCode": status, "headers": {"Content-Type": "text/html; charset=utf-8"}, "body": body}


def _json(data: dict, status: int = 200) -> dict:
    return {
        "statusCode": status,
        "headers": {"Content-Type": "application/json", "Access-Control-Allow-Origin": "*"},
        "body": json.dumps(data, ensure_ascii=False),
    }


def _check_password(body: dict) -> bool:
    if not APP_PASSWORD:
        return True
    return body.get("password") == APP_PASSWORD


def _next_wednesday() -> str:
    today = date.today()
    days = (2 - today.weekday()) % 7
    return (today + timedelta(days=days)).strftime("%Y%m%d")


def _parse_pages(pattern: str, total: int) -> list:
    """페이지 구성 패턴 파싱. update_pptx.parse_pages와 동일."""
    pattern = (pattern or "").strip()
    if not pattern:
        return [total] if total else []

    chunks = [c.strip() for c in re.split(r"[,+]", pattern) if c.strip()]
    is_single_number = len(chunks) == 1 and chunks[0].isdigit()

    result = []
    for chunk in chunks:
        m = re.match(r"^(\d+)\s*[x*X]\s*(\d+)$", chunk)
        if m:
            n, rep = int(m.group(1)), int(m.group(2))
            if n <= 0 or rep <= 0:
                raise ValueError(f"잘못된 패턴: {chunk!r}")
            result.extend([n] * rep)
        elif chunk.isdigit():
            n = int(chunk)
            if n <= 0:
                raise ValueError(f"잘못된 패턴: {chunk!r}")
            result.append(n)
        else:
            raise ValueError(f"잘못된 패턴: {chunk!r}")

    if is_single_number and result:
        per = result[0]
        result = []
        remaining = total
        while remaining > 0:
            take = min(per, remaining)
            result.append(take)
            remaining -= take

    return result


def handle_upload_url(body: dict) -> dict:
    if not _check_password(body):
        return _json({"error": "비밀번호가 틀렸습니다."}, 401)

    filenames = body.get("filenames", [])
    if not filenames:
        return _json({"error": "filenames가 비어있습니다."}, 400)

    session_id = str(uuid.uuid4())[:8]
    urls = []
    for i, name in enumerate(filenames):
        ext = name.rsplit(".", 1)[-1] if "." in name else "jpg"
        key = f"uploads/{session_id}/{i+1}.{ext}"
        url = s3.generate_presigned_url(
            "put_object",
            Params={"Bucket": BUCKET, "Key": key},
            ExpiresIn=300,
        )
        urls.append({"filename": name, "key": key, "url": url})

    return _json({"session_id": session_id, "uploads": urls})


def handle_generate(body: dict) -> dict:
    if not _check_password(body):
        return _json({"error": "비밀번호가 틀렸습니다."}, 401)

    session_id = body.get("session_id")
    keys = body.get("keys", [])
    title = body.get("title", f"{_next_wednesday()} 수요성령집회 콘티")
    pattern = str(body.get("pattern") or body.get("split") or "2")

    if not keys:
        return _json({"error": "keys가 비어있습니다."}, 400)

    try:
        sizes = _parse_pages(pattern, len(keys))
    except ValueError as e:
        return _json({"error": str(e)}, 400)
    if not sizes:
        return _json({"error": "페이지 구성이 비어 있습니다."}, 400)

    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Emu, Pt  # noqa: PLC0415
    from pptx.dml.color import RGBColor  # noqa: PLC0415
    from pptx.enum.shapes import MSO_SHAPE  # noqa: PLC0415
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: PLC0415

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    groups = []
    idx = 0
    for size in sizes:
        if idx >= len(keys):
            break
        groups.append(keys[idx:idx + size])
        idx += size
    if idx < len(keys):
        groups.append(keys[idx:])

    for slide_idx, group in enumerate(groups):
        if not group:
            continue
        slide = prs.slides.add_slide(blank)

        top_offset = Emu(0)
        if slide_idx == 0:
            m = re.match(r"^(\d{4})(\d{2})(\d{2})\s+(.+?)(?:\s+콘티)?$", title)
            if m:
                formatted_date = f"{m.group(1)}.{int(m.group(2))}.{int(m.group(3))}"
                label = m.group(4)
            else:
                formatted_date = ""
                label = title

            NAVY = RGBColor(0x1A, 0x2A, 0x4A)
            WHITE = RGBColor(0xFF, 0xFF, 0xFF)
            box_h = Pt(32)
            label_w = Pt(28 + len(label) * 18)

            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), label_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = NAVY
            box.line.fill.background()
            tf = box.text_frame
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.name = "맑은 고딕"
            run.font.color.rgb = WHITE

            if formatted_date:
                date_box = slide.shapes.add_textbox(label_w + Pt(10), Emu(0), Pt(180), box_h)
                dtf = date_box.text_frame
                dtf.margin_top = Emu(0)
                dtf.margin_bottom = Emu(0)
                dtf.margin_left = Emu(0)
                dtf.margin_right = Emu(0)
                dtf.word_wrap = False
                dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
                dp = dtf.paragraphs[0]
                dp.alignment = PP_ALIGN.LEFT
                drun = dp.add_run()
                drun.text = formatted_date
                drun.font.size = Pt(18)
                drun.font.bold = True
                drun.font.name = "맑은 고딕"
                drun.font.color.rgb = NAVY

            top_offset = box_h + Pt(4)

        n = len(group)
        img_w = slide_w // n
        img_h = slide_h - top_offset

        for i, key in enumerate(group):
            obj = s3.get_object(Bucket=BUCKET, Key=key)
            img_data = io.BytesIO(obj["Body"].read())
            slide.shapes.add_picture(img_data, img_w * i, top_offset, img_w, img_h)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)

    out_key = f"outputs/{session_id}/{title}.pptx"
    s3.put_object(Bucket=BUCKET, Key=out_key, Body=out.read(), ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    return _json({"message": "완료!", "download_key": out_key, "title": title})


def lambda_handler(event, context):
    from pathlib import Path  # noqa: PLC0415
    method = event.get("httpMethod", "GET")
    path = event.get("path", "/")

    if method == "GET" and path == "/":
        _log(event, "page_view")
        html = (Path(__file__).parent / "templates" / "index.html").read_text()
        return _html(html)

    if method == "GET" and path == "/download":
        key = (event.get("queryStringParameters") or {}).get("key", "")
        if not key.startswith("outputs/"):
            return _json({"error": "Invalid key"}, 400)
        filename = key.split("/")[-1]
        url = s3.generate_presigned_url(
            "get_object",
            Params={
                "Bucket": BUCKET,
                "Key": key,
                "ResponseContentDisposition": f"attachment; filename*=UTF-8''{quote(filename)}",
                "ResponseContentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            },
            ExpiresIn=300,
        )
        return {
            "statusCode": 302,
            "headers": {"Location": url},
            "body": "",
        }

    if method == "POST":
        body = json.loads(event.get("body") or "{}")
        if path == "/upload-url":
            return handle_upload_url(body)
        if path == "/log-error":
            _log(event, "client_error", {
                "error": body.get("error", ""),
                "step": body.get("step", ""),
                "ua": (event.get("headers") or {}).get("User-Agent", ""),
            })
            return _json({"ok": True})
        if path == "/generate":
            result = handle_generate(body)
            if result["statusCode"] == 200:
                _log(event, "generate", {
                    "title": body.get("title", ""),
                    "image_count": len(body.get("keys", [])),
                })
            return result

    return _json({"error": "Not found"}, 404)
